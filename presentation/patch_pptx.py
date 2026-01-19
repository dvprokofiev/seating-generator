from pptx import Presentation
from pptx.util import Inches

def patch_video_fullscreen(pptx_path, video_path, slide_index):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    
    # Размеры стандартного широкоэкранного слайда (16:9)
    width = prs.slide_width
    height = prs.slide_height

    # Добавляем видео, растянутое от верхнего левого угла (0,0)
    # до краев слайда
    slide.shapes.add_movie(
        video_path, 
        left=0, 
        top=0, 
        width=width, 
        height=height,
        mime_type='video/mp4'
    )
    
    prs.save('presentation.pptx')

if __name__ == "__main__":
    # Индекс слайда (0 — первый, 3 — четвертый)
    patch_video_fullscreen('dist/presentation.pptx', 'public/demo.mp4', 27)